"""Slide deck generation."""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation

from config import settings
from core.logging import get_logger
from core.models import GeneratedSlide, SlidePlan, TemplateSummary

logger = get_logger(__name__)


class SlideGenerator:
    def __init__(self, template_path: Path, template_summary: TemplateSummary, output_dir: Path) -> None:
        self.template_path = template_path
        self.template_summary = template_summary
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.output_path = self.output_dir / settings.PRESENTATION_NAME
        self.presentation = Presentation(self.template_path)

    def build(self, plans: Iterable[SlidePlan], generated_content: Iterable[GeneratedSlide]) -> Path:
        for plan, content in zip(plans, generated_content):
            layout = self.presentation.slide_layouts[plan.layout.index]
            slide = self.presentation.slides.add_slide(layout)
            self._apply_title(slide, content.title)
            self._apply_body(slide, content.bullet_sentences)
            self._apply_notes(slide, content)
        logger.info("Writing presentation to %s", self.output_path)
        self.presentation.save(self.output_path)
        return self.output_path

    def _apply_title(self, slide, text: str) -> None:
        placeholder = self._get_placeholder(slide, {"title", "centertitle"})
        if placeholder is None:
            logger.debug("No title placeholder found; skipping title")
            return
        placeholder.text = text

    def _apply_body(self, slide, bullets) -> None:
        placeholder = self._get_placeholder(slide, {"body", "content"})
        if placeholder is None:
            logger.debug("No body placeholder found; skipping bullet content")
            return
        text_frame = placeholder.text_frame
        if not bullets:
            text_frame.text = ""
            return
        text_frame.clear()
        first, *rest = bullets
        text_frame.text = first
        for sentence in rest:
            paragraph = text_frame.add_paragraph()
            paragraph.text = sentence
            paragraph.level = 0

    def _apply_notes(self, slide, content: GeneratedSlide) -> None:
        notes = slide.notes_slide
        notes_frame = notes.notes_text_frame
        notes_frame.clear()
        notes_frame.text = f"Outline: {content.outline.title}"
        if content.image_path:
            notes_frame.add_paragraph().text = f"Image: {content.image_path}"
        for key, value in content.notes.items():
            notes_frame.add_paragraph().text = f"{key}: {value}"

    def _get_placeholder(self, slide, accepted_types):
        for placeholder in slide.placeholders:
            placeholder_type = getattr(placeholder.placeholder_format.type, "name", "").lower()
            if placeholder_type in accepted_types:
                return placeholder
        return None


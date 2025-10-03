"""Template analysis utilities."""
from __future__ import annotations

from pathlib import Path
from typing import List

from pptx import Presentation

from config import settings
from core.logging import get_logger
from core.models import TemplateLayout, TemplateSummary

logger = get_logger(__name__)


class TemplateAnalyzer:
    def __init__(self, template_path: Path) -> None:
        self.template_path = template_path

    def analyze(self) -> TemplateSummary:
        logger.info("Analyzing template %s", self.template_path)
        presentation = Presentation(self.template_path)
        layouts: List[TemplateLayout] = []
        for index, layout in enumerate(presentation.slide_layouts):
            name = (layout.name or f"Layout {index}").strip()
            layout_key = name.lower()
            normalized_kind = settings.COMMON_LAYOUT_ALIASES.get(layout_key, "other")
            placeholders = {}
            for placeholder in layout.placeholders:
                fmt = placeholder.placeholder_format
                placeholder_type = getattr(fmt.type, "name", "unknown").lower()
                placeholders[placeholder_type] = fmt.idx
            is_common = normalized_kind in {"title", "content", "content_two"}
            layouts.append(
                TemplateLayout(
                    index=index,
                    name=name,
                    kind=normalized_kind,
                    placeholders=placeholders,
                    is_common=is_common,
                )
            )
        logger.info("Discovered %d layouts", len(layouts))
        return TemplateSummary(layouts=layouts)

